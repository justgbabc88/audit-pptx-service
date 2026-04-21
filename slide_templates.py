"""
Revenue Leak Audit — Slide Templates  (FIXED spacing)
=====================================================
Fixes applied vs. prior version:
  • Titles no longer crowd downstream content when they wrap to 2 lines.
    Each affected slide reserves ~1.6" of vertical space for the title and
    anchors everything below off the same constant (TITLE_BOTTOM).
  • Slide 10 (AI Readiness) key findings block is shortened/capped so it
    does not overflow into the footer — findings are truncated to 2 lines
    of wrap each and the block fits within the slide.
  • System slides (18–20): step cards are taller and text is clamped;
    the right-column first card's metric number no longer rides up on
    top of its label (vertical split fixed).
  • Revenue summary slide: KPI card values shrink if too long and wider
    column for "Overall ROI" so "4.4x return" never clips.
  • Guarantee math box lines now correctly show values in green.
  • Implementation slide phase titles: font auto-shrinks for multi-line
    titles and the items shift down so the first item never collides
    with a wrapped title.
  • Closing slide: safe multi-line rendering (no literal \\n in the text)
    and the urgency line + CTA are on separate y-bands that can't
    overlap even for long copy.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn as _qn
import copy
from lxml import etree

def qn(tag):
    return _qn(tag)

# ─────────────────────────────────────────────
# DESIGN CONSTANTS
# ─────────────────────────────────────────────

# Colors
C_DARK_BG      = RGBColor(0x1C, 0x1B, 0x2E)
C_LIGHT_BG     = RGBColor(0xFF, 0xF8, 0xF0)
C_PURPLE       = RGBColor(0x7B, 0x5E, 0xA7)
C_GREEN        = RGBColor(0x00, 0xC3, 0x89)
C_RED          = RGBColor(0xE6, 0x39, 0x46)
C_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT    = RGBColor(0x1C, 0x1B, 0x2E)
C_MUTED        = RGBColor(0x99, 0x99, 0x99)
C_CARD_BG      = RGBColor(0xFF, 0xFF, 0xFF)
C_HIGHLIGHT    = RGBColor(0xEE, 0xF0, 0xFF)
C_AMBER        = RGBColor(0xE0, 0x7B, 0x00)
C_DARK_CARD    = RGBColor(0x2A, 0x28, 0x45)
C_SCORE_RED    = RGBColor(0xCC, 0x22, 0x22)
C_GREEN_BADGE  = RGBColor(0x00, 0x99, 0x55)
C_AMBER_BADGE  = RGBColor(0xE0, 0x7B, 0x00)
C_RED_BADGE    = RGBColor(0xCC, 0x22, 0x22)
C_TEAL         = RGBColor(0x00, 0xC3, 0x89)
C_FUNNEL_BAR   = RGBColor(0x7B, 0x5E, 0xA7)

# Slide dimensions
W = Inches(13.33)
H = Inches(7.5)

# Watermark
LOGO_PATH = os.path.join(os.path.dirname(__file__), "logo-full-dark.png")
# Native pixel dimensions of logo-full-dark.png (used to preserve aspect ratio)
LOGO_PX_W = 753
LOGO_PX_H = 185

# Margins
LM = Inches(0.6)
RM = Inches(0.6)
TM = Inches(0.5)

# Vertical rhythm — single source of truth for section start after the title
# block. Every light-BG content slide uses these so a wrapping title cannot
# shove content into the next block.
TITLE_TOP       = Inches(0.72)
TITLE_BLOCK_H   = Inches(1.5)   # reserves room for a 2-line title at 36pt
CONTENT_TOP     = TITLE_TOP + TITLE_BLOCK_H  # ≈ 2.22"


# ─────────────────────────────────────────────
# HELPER FUNCTIONS
# ─────────────────────────────────────────────

def new_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def rgb_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def no_fill(shape):
    shape.fill.background()


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        rgb_fill(shape, fill_color)
    else:
        no_fill(shape)
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, radius_pt=6, line_color=None):
    shape = slide.shapes.add_shape(5, left, top, width, height)
    shape.adjustments[0] = radius_pt / 100.0
    if fill_color:
        rgb_fill(shape, fill_color)
    else:
        no_fill(shape)
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _clean_text(text):
    """Convert any literal '\\n' sequences into real newlines. Defensive
    against data that has been JSON-encoded or copy-pasted with backslashes.
    """
    if text is None:
        return ""
    if not isinstance(text, str):
        text = str(text)
    # Replace literal backslash-n with real newline
    return text.replace("\\n", "\n")


def add_textbox(slide, left, top, width, height, text, font_name, font_size,
                bold=False, italic=False, color=C_DARK_TEXT, align=PP_ALIGN.LEFT,
                word_wrap=True, line_spacing=None):
    text = _clean_text(text)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if line_spacing:
            pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, _qn('a:lnSpc'))
            spcPts = etree.SubElement(lnSpc, _qn('a:spcPts'))
            spcPts.set('val', str(int(line_spacing * 100)))

    txBody = txBox.text_frame._txBody
    bodyPr = txBody.find(_qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('lIns', '0')
        bodyPr.set('rIns', '0')
        bodyPr.set('tIns', '0')
        bodyPr.set('bIns', '0')
    return txBox


def add_vcenter_text(slide, left, top, width, height, text, font_name, font_size,
                     bold=False, italic=False, color=C_DARK_TEXT, align=PP_ALIGN.CENTER):
    """Vertically centred text using a real textbox with MSO_ANCHOR.MIDDLE.
    Previously this used a manual top-pad calculation that drifted as text
    wrapped; using the built-in vertical anchor is both simpler and reliable.
    """
    text = _clean_text(text)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # pure black for legibility

    txBody = txBox.text_frame._txBody
    bodyPr = txBody.find(_qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('lIns', '0')
        bodyPr.set('rIns', '0')
        bodyPr.set('tIns', '0')
        bodyPr.set('bIns', '0')
    return txBox


def add_accent_bar(slide, dark=True):
    height = Inches(0.12) if dark else Inches(0.08)
    bar = add_rect(slide, 0, 0, W, height, fill_color=C_PURPLE)
    return bar


def add_footer(slide, brand_name, slide_number=None):
    footer_y = H - Inches(0.35)
    footer_h = Inches(0.25)
    add_textbox(slide, LM, footer_y, Inches(6), footer_h,
                f"{brand_name}  \u00b7  Revenue Leak Audit  \u00b7  Confidential",
                "Calibri", 10, color=C_MUTED)


def add_section_label(slide, text, left=None, top=Inches(0.45), dark=False):
    if left is None:
        left = LM
    color = C_PURPLE
    add_textbox(slide, left, top, W - left - RM, Inches(0.25),
                text, "Arial", 11, bold=True, color=color)


def _auto_title_size(text, default=36, base_chars=46, min_size=24):
    """Shrink title font size if the title is long so it fits on 1-2 lines
    rather than pushing into the content below."""
    text = _clean_text(text)
    # Rough heuristic: if >base_chars characters, scale down proportionally
    if len(text) <= base_chars:
        return default
    scale = base_chars / max(1, len(text))
    return max(min_size, int(default * scale))


def add_title(slide, text, top=TITLE_TOP, font_size=40, color=C_DARK_TEXT,
              width=None, height=None, auto_shrink=False):
    """Render a slide title inside a fixed-height block (default TITLE_BLOCK_H).
    If auto_shrink=True, font size is reduced for long titles so the title
    fits in the reserved block without overflowing into content below.
    """
    if width is None:
        width = W - LM - RM
    if height is None:
        height = TITLE_BLOCK_H
    if auto_shrink:
        font_size = _auto_title_size(text, default=font_size)
    add_textbox(slide, LM, top, width, height,
                text, "Georgia", font_size, bold=True, color=color)


def set_shape_text(shape, text, font_name, font_size, bold=False, italic=False,
                   color=C_DARK_TEXT, align=PP_ALIGN.LEFT):
    text = _clean_text(text)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_multi_para_textbox(slide, left, top, width, height, paragraphs, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('lIns', '0'); bodyPr.set('rIns', '0')
        bodyPr.set('tIns', '0'); bodyPr.set('bIns', '0')

    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = para.get('align', PP_ALIGN.LEFT)
        if para.get('space_before'):
            pPr = p._p.get_or_add_pPr()
            spcBef = etree.SubElement(pPr, qn('a:spcBef'))
            spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
            spcPts.set('val', str(int(para['space_before'] * 100)))
        run = p.add_run()
        run.text = _clean_text(para.get('text', ''))
        run.font.name = para.get('font_name', 'Calibri')
        run.font.size = Pt(para.get('font_size', 14))
        run.font.bold = para.get('bold', False)
        run.font.italic = para.get('italic', False)
        run.font.color.rgb = para.get('color', C_DARK_TEXT)
    return txBox


def add_circle(slide, cx, cy, diameter, fill_color, text, text_color=C_WHITE,
               font_size=14, bold=True):
    r = diameter / 2
    shape = slide.shapes.add_shape(9, int(cx - r), int(cy - r), int(diameter), int(diameter))
    rgb_fill(shape, fill_color)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Georgia"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    tf.auto_size = None
    tf.word_wrap = False
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_logo_watermark(slide, image_path=LOGO_PATH, width_in=1.8):
    """Place the logo-full-dark.png mark in the bottom-right corner of the
    slide at full opacity, where the slide number used to live."""
    if not os.path.exists(image_path):
        return None

    wm_w = Inches(width_in)
    wm_h = Emu(int(wm_w * LOGO_PX_H / LOGO_PX_W))
    left = W - wm_w - Inches(0.4)
    top  = H - wm_h - Inches(0.15)

    return slide.shapes.add_picture(image_path, left, top, width=wm_w, height=wm_h)


# ─────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────

def slide_01_cover(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_DARK_BG)
    add_accent_bar(slide, dark=True)

    add_textbox(slide, LM, Inches(0.35), Inches(5), Inches(0.3),
                "REVENUE LEAK AUDIT", "Arial", 11, bold=True, color=C_PURPLE)

    add_textbox(slide, LM, Inches(1.3), W - LM - RM, Inches(2.1),
                d["client_name"], "Georgia", 72, bold=True, color=C_WHITE)

    add_textbox(slide, LM, Inches(3.45), Inches(9), Inches(0.38),
                "AI Readiness Assessment & Revenue Recovery Analysis",
                "Calibri", 18, color=C_MUTED)

    add_rect(slide, LM, Inches(3.98), Inches(1.2), Inches(0.05), fill_color=C_RED)

    add_textbox(slide, LM, Inches(4.22), Inches(8), Inches(0.3),
                f"Prepared for {d['prepared_for']} \u2014 {d['client_name']}",
                "Calibri", 14, color=C_WHITE)
    add_textbox(slide, LM, Inches(4.56), Inches(4), Inches(0.3),
                d["date"], "Calibri", 14, color=C_WHITE)

    add_textbox(slide, LM, Inches(5.9), Inches(3), Inches(0.3),
                "CONFIDENTIAL", "Arial", 11, bold=True, color=C_RED)


# ─────────────────────────────────────────────
# SLIDE 2 — AGENDA
# ─────────────────────────────────────────────

def slide_02_agenda(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_section_label(slide, "TODAY'S AGENDA")
    add_title(slide, "What we'll walk through together", top=TITLE_TOP, font_size=40)

    items = [
        ("01", "Your Industry",          d["agenda_01_desc"]),
        ("02", "Your Business Today",    "What we found when we audited your " + d["term_business"]),
        ("03", "Your AI Readiness Score","6 dimensions, scored with specific findings"),
        ("04", "The Revenue Leak",       "Exact dollars you're leaving on the table"),
        ("05", "Your AI Systems",        "The specific bots we recommend and how each works"),
        ("06", "The Guarantee",          "Our 90-day ROI guarantee \u2014 risk-free"),
        ("07", "Next Steps",             "Select your bots and we start building"),
    ]

    row_h = Inches(0.62)
    row_top = Inches(2.3)
    row_w = W - LM - RM

    for i, (num, title, desc) in enumerate(items):
        bg = C_HIGHLIGHT if i % 2 == 0 else C_WHITE
        add_rounded_rect(slide, LM, row_top + i * row_h, row_w, row_h - Inches(0.04),
                         fill_color=bg, radius_pt=4)
        add_textbox(slide, LM + Inches(0.12), row_top + i * row_h + Inches(0.12),
                    Inches(0.5), Inches(0.38),
                    num, "Georgia", 22, bold=True, color=C_PURPLE)
        add_textbox(slide, LM + Inches(0.65), row_top + i * row_h + Inches(0.14),
                    Inches(2.8), Inches(0.34),
                    title, "Calibri", 15, bold=True, color=C_DARK_TEXT)
        add_textbox(slide, LM + Inches(3.6), row_top + i * row_h + Inches(0.16),
                    Inches(8.5), Inches(0.3),
                    desc, "Calibri", 14, color=C_MUTED)

    add_footer(slide, d["brand_name"], 2)


# ─────────────────────────────────────────────
# SLIDE 3 — PART ONE DIVIDER
# ─────────────────────────────────────────────

def slide_03_part_one(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_DARK_BG)
    add_accent_bar(slide, dark=True)

    add_rect(slide, 0, 0, Inches(0.08), H, fill_color=C_PURPLE)

    add_textbox(slide, LM, Inches(1.8), Inches(4), Inches(0.3),
                "PART ONE", "Arial", 11, bold=True, color=C_PURPLE)

    add_textbox(slide, LM, Inches(2.25), W - LM - RM, Inches(2.0),
                d["part_one_title"], "Georgia", 52, bold=True, color=C_WHITE)

    add_textbox(slide, LM, Inches(4.45), Inches(9), Inches(0.4),
                d["part_one_subtitle"], "Calibri", 18, color=C_MUTED)


# ─────────────────────────────────────────────
# SLIDE 4 — CLIENT JOURNEY
# ─────────────────────────────────────────────

def slide_04_journey(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_section_label(slide, d["journey_label"])
    # Auto-shrink so long journey titles don't wrap to 3 lines
    add_title(slide, d["journey_title"], top=TITLE_TOP, font_size=32, auto_shrink=True)

    stages = d["journey_stages"]
    n = len(stages)
    total_w = W - LM - RM
    card_w = (total_w - Inches(0.1) * (n - 1)) / n
    card_h = Inches(1.42)
    # Anchor cards OFF the reserved title block, not a hard-coded y.
    card_top = CONTENT_TOP + Inches(0.1)  # ≈ 2.32"

    stage_colors = [
        RGBColor(0x5B, 0x8D, 0xE8), RGBColor(0x7B, 0x5E, 0xA7),
        RGBColor(0xE0, 0x7B, 0x00), RGBColor(0x00, 0xC3, 0x89),
        RGBColor(0x2E, 0x8B, 0x6E), RGBColor(0x00, 0x99, 0xCC),
    ]

    for i, stage in enumerate(stages):
        cx = LM + i * (card_w + Inches(0.1))
        add_rounded_rect(slide, cx, card_top, card_w, card_h,
                         fill_color=C_CARD_BG, radius_pt=4,
                         line_color=RGBColor(0xDD, 0xDD, 0xDD))
        color = stage_colors[i % len(stage_colors)]
        add_rect(slide, cx, card_top, card_w, Inches(0.06), fill_color=color)
        add_textbox(slide, cx + Inches(0.08), card_top + Inches(0.1),
                    card_w - Inches(0.16), Inches(0.28),
                    stage["name"], "Arial", 9, bold=True, color=color,
                    align=PP_ALIGN.CENTER)
        add_vcenter_text(slide, cx + Inches(0.05), card_top + Inches(0.42),
                         card_w - Inches(0.1), card_h - Inches(0.42),
                         stage["desc"], "Calibri", 11,
                         color=RGBColor(0x55, 0x55, 0x55))

        if stage.get("pct_lost"):
            add_textbox(slide, cx, card_top + card_h + Inches(0.1),
                        card_w, Inches(0.28),
                        stage["pct_lost"] + " lost", "Calibri", 13, bold=True,
                        color=C_RED, align=PP_ALIGN.CENTER)
            add_textbox(slide, cx, card_top + card_h + Inches(0.38),
                        card_w, Inches(0.24),
                        stage.get("reason", ""), "Calibri", 10,
                        color=C_MUTED, align=PP_ALIGN.CENTER)

    add_textbox(slide, LM, card_top + card_h + Inches(0.78), W - LM - RM, Inches(0.35),
                d["journey_closing"], "Calibri", 13, italic=True, color=C_RED)

    add_footer(slide, d["brand_name"], 4)


# ─────────────────────────────────────────────
# SLIDE 5 — INDUSTRY STATS
# ─────────────────────────────────────────────

def slide_05_stats(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_section_label(slide, d["stat_section_label"])
    add_title(slide, "What the data says about businesses like yours",
              top=TITLE_TOP, font_size=32, auto_shrink=True)

    stats = d["industry_stats"]
    card_w = Inches(5.8)
    card_h = Inches(1.95)
    gap = Inches(0.25)
    top_row_y = CONTENT_TOP + Inches(0.1)   # anchored off reserved title block
    bot_row_y = top_row_y + card_h + Inches(0.2)
    positions = [
        (LM, top_row_y),
        (LM + card_w + gap, top_row_y),
        (LM, bot_row_y),
        (LM + card_w + gap, bot_row_y),
    ]

    for i, (stat, (cx, cy)) in enumerate(zip(stats, positions)):
        add_rounded_rect(slide, cx, cy, card_w, card_h,
                         fill_color=C_CARD_BG, radius_pt=4,
                         line_color=RGBColor(0xE0, 0xE0, 0xE0))
        add_textbox(slide, cx + Inches(0.18), cy + Inches(0.1),
                    card_w - Inches(0.36), Inches(0.78),
                    stat["number"], "Georgia", 44, bold=True, color=C_DARK_TEXT)
        add_textbox(slide, cx + Inches(0.18), cy + Inches(0.92),
                    card_w - Inches(0.36), Inches(0.62),
                    stat["label"], "Calibri", 12, color=C_DARK_TEXT)
        add_textbox(slide, cx + Inches(0.18), cy + Inches(1.58),
                    card_w - Inches(0.36), Inches(0.26),
                    stat["source"], "Calibri", 10, italic=True, color=C_MUTED)

    add_footer(slide, d["brand_name"], 5)


# ─────────────────────────────────────────────
# SLIDE 6 — LEAD LIFECYCLE
# ─────────────────────────────────────────────

def slide_06_lifecycle(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_section_label(slide, "THE LEAD LIFECYCLE")
    # Shrink title so it sits in one line when possible and leaves room for intro
    add_title(slide, d["lifecycle_title"], top=TITLE_TOP, font_size=28, auto_shrink=True)

    # Intro paragraph — anchored to reserved title block, not a hard-coded y
    intro_top = CONTENT_TOP
    add_textbox(slide, LM, intro_top, W - LM - RM, Inches(0.5),
                d["lifecycle_intro"], "Calibri", 13, color=C_DARK_TEXT)

    stages = d["lifecycle_stages"]
    n = 6
    total_w = W - LM - RM
    card_w = (total_w - Inches(0.08) * (n - 1)) / n
    card_h = Inches(1.95)
    card_top = intro_top + Inches(0.7)

    for i, stage in enumerate(stages):
        cx = LM + i * (card_w + Inches(0.08))
        add_rounded_rect(slide, cx, card_top, card_w, card_h,
                         fill_color=C_CARD_BG, radius_pt=4,
                         line_color=RGBColor(0xDD, 0xDD, 0xDD))
        add_textbox(slide, cx + Inches(0.06), card_top + Inches(0.1),
                    card_w - Inches(0.12), Inches(0.32),
                    stage["name"], "Arial", 9, bold=True, color=C_PURPLE,
                    align=PP_ALIGN.CENTER)
        add_vcenter_text(slide, cx + Inches(0.05), card_top + Inches(0.44),
                         card_w - Inches(0.1), card_h - Inches(0.44) - Inches(0.44),
                         stage["desc"], "Calibri", 10,
                         color=RGBColor(0x55, 0x55, 0x55))

        status = stage["status"].upper()
        badge_color = C_GREEN_BADGE if status == "WORKING" else (
            C_AMBER_BADGE if status == "MANUAL" else C_RED_BADGE)
        badge_h = Inches(0.32)
        badge_top = card_top + card_h - badge_h - Inches(0.1)
        badge = add_rounded_rect(slide, cx + Inches(0.1), badge_top,
                                  card_w - Inches(0.2), badge_h,
                                  fill_color=badge_color, radius_pt=4)
        set_shape_text(badge, status, "Arial", 9, bold=True,
                       color=C_WHITE, align=PP_ALIGN.CENTER)
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, LM, card_top + card_h + Inches(0.18), W - LM - RM, Inches(0.45),
                d["lifecycle_summary"], "Calibri", 13, italic=True, color=C_RED)

    add_footer(slide, d["brand_name"], 6)


# ─────────────────────────────────────────────
# SLIDE 7 — REVENUE FUNNEL
# ─────────────────────────────────────────────

def slide_07_funnel(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_section_label(slide, "YOUR REVENUE FUNNEL")
    add_title(slide, "Where your monthly leads actually end up",
              top=TITLE_TOP, font_size=32, auto_shrink=True)

    funnel = d["funnel_rows"]
    bar_top = CONTENT_TOP
    bar_h = Inches(0.56)
    bar_gap = Inches(0.06)
    max_w = W - LM - RM - Inches(1.8)

    widths = [1.0, 0.82, 0.62, 0.42, 0.28]

    for i, (row, w_frac) in enumerate(zip(funnel, widths)):
        bar_w = max_w * w_frac
        y = bar_top + i * (bar_h + bar_gap)
        add_rounded_rect(slide, LM, y, bar_w, bar_h,
                         fill_color=C_FUNNEL_BAR, radius_pt=3)
        add_textbox(slide, LM + Inches(0.15), y + Inches(0.1),
                    bar_w - Inches(1.4), Inches(0.36),
                    row["label"], "Calibri", 13, color=C_WHITE)
        add_textbox(slide, LM + bar_w - Inches(1.2), y + Inches(0.05),
                    Inches(1.1), Inches(0.46),
                    str(row["value"]), "Georgia", 26, bold=True,
                    color=C_WHITE, align=PP_ALIGN.RIGHT)
        if row.get("lost"):
            add_textbox(slide, LM + bar_w + Inches(0.15), y + Inches(0.1),
                        Inches(1.5), Inches(0.36),
                        row["lost"], "Calibri", 13, bold=True, color=C_RED)

    # Insight paragraph — anchored AFTER the last bar so it can't collide
    insight_top = bar_top + 5 * (bar_h + bar_gap) + Inches(0.2)
    add_textbox(slide, LM, insight_top, W - LM - RM, Inches(0.8),
                d["funnel_insight"], "Calibri", 13, italic=True, color=C_DARK_TEXT)

    add_footer(slide, d["brand_name"], 7)


# ─────────────────────────────────────────────
# SLIDE 8 — PART TWO DIVIDER
# ─────────────────────────────────────────────

def slide_08_part_two(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_DARK_BG)
    add_accent_bar(slide, dark=True)
    add_rect(slide, 0, 0, Inches(0.08), H, fill_color=C_PURPLE)

    add_textbox(slide, LM, Inches(1.8), Inches(4), Inches(0.3),
                "PART TWO", "Arial", 11, bold=True, color=C_PURPLE)
    add_textbox(slide, LM, Inches(2.25), W - LM - RM, Inches(1.2),
                "Your Business Today", "Georgia", 52, bold=True, color=C_WHITE)
    add_textbox(slide, LM, Inches(3.6), Inches(9), Inches(0.4),
                f"What we found when we audited {d['client_name']}",
                "Calibri", 18, color=C_MUTED)


# ─────────────────────────────────────────────
# SLIDE 9 — YOUR NUMBERS
# ─────────────────────────────────────────────

def slide_09_numbers(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_section_label(slide, "YOUR NUMBERS")
    add_title(slide, f"{d['client_name']} \u2014 current state",
              top=TITLE_TOP, font_size=32, auto_shrink=True)

    cards = d["number_cards"]
    card_w = Inches(3.8)
    card_h = Inches(2.1)
    gap_x = Inches(0.28)
    gap_y = Inches(0.28)
    top_y = CONTENT_TOP

    positions = [
        (LM,                        top_y),
        (LM + card_w + gap_x,       top_y),
        (LM + 2*(card_w + gap_x),   top_y),
        (LM,                        top_y + card_h + gap_y),
        (LM + card_w + gap_x,       top_y + card_h + gap_y),
        (LM + 2*(card_w + gap_x),   top_y + card_h + gap_y),
    ]

    for card, (cx, cy) in zip(cards, positions):
        add_rounded_rect(slide, cx, cy, card_w, card_h,
                         fill_color=C_CARD_BG, radius_pt=4,
                         line_color=RGBColor(0xE0, 0xE0, 0xE0))
        # Auto-shrink big value for long strings like "$12,000"
        value_str = str(card["value"])
        value_size = 44 if len(value_str) <= 6 else (36 if len(value_str) <= 8 else 30)
        add_textbox(slide, cx + Inches(0.18), cy + Inches(0.14),
                    card_w - Inches(0.36), Inches(0.9),
                    value_str, "Georgia", value_size, bold=True, color=C_DARK_TEXT)
        add_textbox(slide, cx + Inches(0.18), cy + Inches(1.0),
                    card_w - Inches(0.36), Inches(0.55),
                    card["label"], "Calibri", 14, color=C_DARK_TEXT)
        add_textbox(slide, cx + Inches(0.18), cy + Inches(1.58),
                    card_w - Inches(0.36), Inches(0.38),
                    card["subnote"], "Calibri", 11, italic=True, color=C_MUTED)

    add_footer(slide, d["brand_name"], 9)


# ─────────────────────────────────────────────
# SLIDE 10 — AUDIT METHODOLOGY
# ─────────────────────────────────────────────

def slide_10_methodology(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_section_label(slide, "AUDIT METHODOLOGY")
    add_title(slide, "What we actually tested", top=TITLE_TOP, font_size=32, auto_shrink=True)

    add_textbox(slide, LM, CONTENT_TOP, W - LM - RM, Inches(0.35),
                f"We didn't just ask questions \u2014 we tested your {d['term_business']} "
                f"as a real {d['term_client']} would.",
                "Calibri", 13, italic=True, color=C_DARK_TEXT)

    tests = d["audit_tests"]
    col_w = Inches(5.9)
    row_h = Inches(1.3)
    gap_x = Inches(0.26)
    gap_y = Inches(0.1)
    start_y = CONTENT_TOP + Inches(0.45)

    for i, test in enumerate(tests):
        col = i % 2
        row = i // 2
        cx = LM + col * (col_w + gap_x)
        cy = start_y + row * (row_h + gap_y)

        add_rounded_rect(slide, cx, cy, col_w, row_h,
                         fill_color=C_CARD_BG, radius_pt=4,
                         line_color=RGBColor(0xE0, 0xE0, 0xE0))

        add_circle(slide, cx + Inches(0.35), cy + row_h / 2,
                   Inches(0.42), C_PURPLE,
                   str(i + 1), font_size=16)

        add_textbox(slide, cx + Inches(0.7), cy + Inches(0.1),
                    col_w - Inches(0.85), Inches(0.38),
                    test["name"], "Calibri", 14, bold=True, color=C_DARK_TEXT)
        add_textbox(slide, cx + Inches(0.7), cy + Inches(0.5),
                    col_w - Inches(0.85), Inches(0.65),
                    test["finding"], "Calibri", 12, color=C_MUTED)

    add_footer(slide, d["brand_name"], 10)


# ─────────────────────────────────────────────
# SLIDE 11 — AI READINESS SCORE
# ─────────────────────────────────────────────

def slide_11_score(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_DARK_BG)
    add_accent_bar(slide, dark=True)

    add_textbox(slide, LM, Inches(0.3), Inches(8), Inches(0.3),
                "YOUR AI READINESS SCORE", "Arial", 11, bold=True, color=C_PURPLE)

    dims = d["dimensions"]
    bar_top = Inches(0.78)
    bar_h = Inches(0.48)
    bar_gap = Inches(0.08)
    bar_track_w = Inches(4.5)
    label_w = Inches(2.4)

    for i, dim in enumerate(dims):
        y = bar_top + i * (bar_h + bar_gap)
        score_frac = dim["score"] / 100.0

        add_textbox(slide, LM, y + Inches(0.06), label_w, Inches(0.36),
                    dim["name"], "Calibri", 13, color=C_WHITE)
        add_textbox(slide, LM + label_w, y + Inches(0.08), Inches(0.6), Inches(0.3),
                    dim["weight"], "Calibri", 11, color=C_MUTED)

        track_x = LM + label_w + Inches(0.7)
        add_rect(slide, track_x, y + Inches(0.12),
                 bar_track_w, Inches(0.24),
                 fill_color=RGBColor(0x3A, 0x38, 0x5A))

        score = dim["score"]
        fill_color = C_RED if score < 30 else (C_AMBER if score < 60 else C_GREEN)
        fill_w = max(Inches(0.15), bar_track_w * score_frac)
        add_rect(slide, track_x, y + Inches(0.12),
                 fill_w, Inches(0.24), fill_color=fill_color)

        add_textbox(slide, track_x + bar_track_w + Inches(0.1), y + Inches(0.04),
                    Inches(0.5), Inches(0.36),
                    str(score), "Calibri", 16, bold=True, color=fill_color)

    # KEY FINDINGS — tighter spacing AND room-aware truncation so the block
    # fits comfortably above the footer no matter how many findings are passed.
    kf_y = bar_top + 6 * (bar_h + bar_gap) + Inches(0.15)
    add_textbox(slide, LM, kf_y, Inches(5), Inches(0.28),
                "KEY FINDINGS", "Arial", 10, bold=True, color=C_PURPLE)

    findings = d.get("key_findings", [])
    # Available height between kf_y + 0.34" (first line) and footer at H - 0.45"
    footer_safe = H - Inches(0.5)
    available_h = footer_safe - (kf_y + Inches(0.34))
    # Allow ~0.32" per finding (single line) — cap list length to what fits.
    line_h = Inches(0.32)
    max_findings = int(available_h / line_h)
    findings = findings[:max_findings]
    # Compact width so each finding stays on ONE line (prevents overflow).
    for i, finding in enumerate(findings):
        # Trim text that would wrap — keep it readable
        compact = _shorten_for_single_line(finding, max_chars=110)
        add_textbox(slide, LM, kf_y + Inches(0.34) + i * line_h,
                    Inches(7.9), line_h,
                    "\u2192  " + compact, "Calibri", 11, color=C_WHITE)

    # Score circle (right side) — moved up and made slightly smaller so the
    # /100 label and status badge below it never overlap the circle.
    circle_cx = Inches(10.8)
    circle_cy = Inches(2.4)
    circle_d = Inches(2.4)
    add_circle(slide, circle_cx, circle_cy, circle_d,
               C_SCORE_RED, str(d["overall_score"]),
               font_size=64)

    # /100 — sits well below the circle (circle bottom is at cy + r = 3.6")
    add_textbox(slide, circle_cx - Inches(1.4), Inches(3.75),
                Inches(2.8), Inches(0.4),
                "/100", "Calibri", 22, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Status badge — further down, in its own y-band with no chance of
    # overlapping the circle even for 2-line labels.
    add_textbox(slide, circle_cx - Inches(1.9), Inches(4.3),
                Inches(3.8), Inches(0.8),
                d["score_status_label"], "Arial", 13, bold=True,
                color=C_RED, align=PP_ALIGN.CENTER)

    add_footer(slide, d["brand_name"], 11)


def _shorten_for_single_line(text, max_chars=110):
    text = _clean_text(text).replace("\n", " ")
    # collapse whitespace
    text = " ".join(text.split())
    if len(text) <= max_chars:
        return text
    # Cut at last space before limit, add ellipsis
    cut = text[:max_chars]
    last_space = cut.rfind(" ")
    if last_space > 60:
        cut = cut[:last_space]
    return cut.rstrip(",.;:\u2014- ") + "\u2026"


# ─────────────────────────────────────────────
# DIMENSION DEEP DIVE
# ─────────────────────────────────────────────

def slide_deep_dive(prs, d, dive, slide_number):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_section_label(slide, f"DIMENSION DEEP DIVE \u2014 {dive['score']}/100")

    badge = add_rounded_rect(slide, W - RM - Inches(2.6), Inches(0.3),
                              Inches(2.6), Inches(1.1),
                              fill_color=C_SCORE_RED, radius_pt=4)
    add_textbox(slide, W - RM - Inches(2.55), Inches(0.32),
                Inches(2.5), Inches(0.62),
                f"{dive['score']}/100", "Georgia", 36, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, W - RM - Inches(2.55), Inches(0.92),
                Inches(2.5), Inches(0.3),
                dive["subtitle"], "Calibri", 11,
                color=C_WHITE, align=PP_ALIGN.CENTER)

    add_title(slide, dive["name"], top=TITLE_TOP, font_size=32,
              width=W - LM - RM - Inches(3.0), auto_shrink=True)

    subs = dive["sub_scores"]
    sub_card_w = Inches(5.7)
    sub_card_h = Inches(0.72)
    sub_gap_x = Inches(0.28)
    sub_gap_y = Inches(0.1)
    sub_top = Inches(1.75)

    for i, sub in enumerate(subs):
        col = i % 2
        row = i // 2
        cx = LM + col * (sub_card_w + sub_gap_x)
        cy = sub_top + row * (sub_card_h + sub_gap_y)

        add_rounded_rect(slide, cx, cy, sub_card_w, sub_card_h,
                         fill_color=C_CARD_BG, radius_pt=4,
                         line_color=RGBColor(0xE0, 0xE0, 0xE0))

        add_textbox(slide, cx + Inches(0.12), cy + Inches(0.05),
                    Inches(0.6), Inches(0.62),
                    str(sub["score"]), "Georgia", 28, bold=True, color=C_PURPLE)
        add_textbox(slide, cx + Inches(0.75), cy + Inches(0.06),
                    sub_card_w - Inches(0.9), Inches(0.3),
                    sub["label"], "Calibri", 13, bold=True, color=C_DARK_TEXT)
        add_textbox(slide, cx + Inches(0.75), cy + Inches(0.36),
                    sub_card_w - Inches(0.9), Inches(0.28),
                    sub["note"], "Calibri", 11, color=C_MUTED)

    af_top = Inches(3.42)
    af_h = Inches(1.3)
    add_rect(slide, LM, af_top, W - LM - RM, af_h,
             fill_color=RGBColor(0xFF, 0xF0, 0xF0))
    add_rect(slide, LM, af_top, Inches(0.05), af_h, fill_color=C_RED)
    add_textbox(slide, LM + Inches(0.15), af_top + Inches(0.08),
                Inches(3), Inches(0.28),
                "AUDIT FINDING", "Arial", 10, bold=True, color=C_RED)
    add_textbox(slide, LM + Inches(0.15), af_top + Inches(0.36),
                W - LM - RM - Inches(0.3), Inches(0.85),
                dive["audit_finding"], "Calibri", 12, color=C_DARK_TEXT)

    ri_top = af_top + af_h + Inches(0.12)
    add_rect(slide, LM, ri_top, W - LM - RM, Inches(1.0),
             fill_color=RGBColor(0xF0, 0xFF, 0xF6))
    add_rect(slide, LM, ri_top, Inches(0.05), Inches(1.0), fill_color=C_GREEN)
    add_textbox(slide, LM + Inches(0.15), ri_top + Inches(0.06),
                Inches(3), Inches(0.28),
                "REVENUE IMPACT", "Arial", 10, bold=True, color=C_GREEN)
    add_textbox(slide, LM + Inches(0.15), ri_top + Inches(0.34),
                W - LM - RM - Inches(0.3), Inches(0.58),
                dive["revenue_impact"], "Calibri", 12, color=C_DARK_TEXT)

    add_footer(slide, d["brand_name"], slide_number)


# ─────────────────────────────────────────────
# SLIDE 15 — BOTTOM LINE
# ─────────────────────────────────────────────

def slide_15_bottom_line(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_DARK_BG)
    add_accent_bar(slide, dark=True)

    add_textbox(slide, LM, Inches(0.35), Inches(5), Inches(0.28),
                "THE BOTTOM LINE", "Arial", 11, bold=True, color=C_PURPLE)

    add_textbox(slide, LM, Inches(0.85), W - LM - RM, Inches(0.55),
                f"Your {d['bottom_line_biz_word']} is leaving",
                "Calibri", 30, color=C_WHITE)

    add_textbox(slide, LM, Inches(1.38), W - LM - RM, Inches(1.85),
                d["annual_leak"], "Georgia", 96, bold=True, color=C_GREEN)

    add_textbox(slide, LM, Inches(3.28), W - LM - RM, Inches(0.55),
                "per year on the table.", "Calibri", 30, color=C_WHITE)

    add_textbox(slide, LM, Inches(4.05), W - LM - RM, Inches(0.75),
                d["bottom_line_context"], "Calibri", 15, color=C_MUTED)

    add_textbox(slide, LM, Inches(5.05), W - LM - RM, Inches(1.0),
                d["closing_line"], "Georgia", 19, italic=True, color=C_WHITE)


# ─────────────────────────────────────────────
# SLIDE 16 — PART THREE DIVIDER
# ─────────────────────────────────────────────

def slide_16_part_three(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_DARK_BG)
    add_accent_bar(slide, dark=True)
    add_rect(slide, 0, 0, Inches(0.08), H, fill_color=C_PURPLE)

    add_textbox(slide, LM, Inches(1.8), Inches(4), Inches(0.3),
                "PART THREE", "Arial", 11, bold=True, color=C_PURPLE)
    add_textbox(slide, LM, Inches(2.25), W - LM - RM, Inches(1.5),
                "Your AI Revenue\nRecovery Systems", "Georgia", 52, bold=True,
                color=C_WHITE)
    add_textbox(slide, LM, Inches(4.05), Inches(9), Inches(0.4),
                f"Recommended systems custom-built for {d['client_name']}",
                "Calibri", 18, color=C_MUTED)


# ─────────────────────────────────────────────
# SLIDE 17 — HOW IT WORKS
# ─────────────────────────────────────────────

def slide_17_how_it_works(prs, d):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_section_label(slide, "HOW IT WORKS")
    add_title(slide, f"What changes in your {d['term_client']} journey",
              top=TITLE_TOP, font_size=32, auto_shrink=True)

    before_steps = [
        "Lead calls in",
        "Voicemail / wait",
        "One email reply",
        "No follow-up",
        d["before_dropoff_pct"],
        "Never return",
    ]

    col_w = (W - LM - RM) / 6 - Inches(0.07)
    step_h = Inches(1.3)

    add_textbox(slide, LM, CONTENT_TOP - Inches(0.05), Inches(6), Inches(0.3),
                "BEFORE (current)", "Arial", 11, bold=True, color=C_RED)
    step_top_before = CONTENT_TOP + Inches(0.3)

    for i, step in enumerate(before_steps):
        cx = LM + i * (col_w + Inches(0.07))
        add_rounded_rect(slide, cx, step_top_before, col_w, step_h,
                         fill_color=RGBColor(0xF0, 0xEE, 0xEE), radius_pt=3)
        add_vcenter_text(slide, cx + Inches(0.05), step_top_before,
                         col_w - Inches(0.1), step_h,
                         step, "Calibri", 12, color=C_DARK_TEXT)
        if i < 5:
            add_textbox(slide, cx + col_w, step_top_before + Inches(0.48),
                        Inches(0.07), Inches(0.3),
                        "\u2014", "Calibri", 11, color=C_MUTED, align=PP_ALIGN.CENTER)

    after_label_top = step_top_before + step_h + Inches(0.22)
    add_textbox(slide, LM, after_label_top, Inches(6), Inches(0.3),
                "AFTER (with AI systems)", "Arial", 11, bold=True, color=C_GREEN)
    step_top_after = after_label_top + Inches(0.35)

    after_steps = [
        "Lead calls in",
        "AI answers in\n<3 rings, 24/7",
        "7-touch nurture\nfor non-bookers",
        "Smart reminders\nslash no-shows",
        d["after_step_5"],
        "Reviews, referrals,\nreactivation",
    ]

    for i, step in enumerate(after_steps):
        cx = LM + i * (col_w + Inches(0.07))
        add_rounded_rect(slide, cx, step_top_after, col_w, step_h,
                         fill_color=RGBColor(0xE8, 0xF8, 0xF2), radius_pt=3)
        add_vcenter_text(slide, cx + Inches(0.05), step_top_after,
                         col_w - Inches(0.1), step_h,
                         step, "Calibri", 11, color=C_DARK_TEXT)
        if i < 5:
            add_textbox(slide, cx + col_w, step_top_after + Inches(0.48),
                        Inches(0.07), Inches(0.3),
                        "\u2014", "Calibri", 11, color=C_MUTED, align=PP_ALIGN.CENTER)

    add_textbox(slide, LM, step_top_after + step_h + Inches(0.15), W - LM - RM, Inches(0.45),
                d["how_it_works_closing"], "Calibri", 13, italic=True, color=C_GREEN)

    add_footer(slide, d["brand_name"], 17)


# ─────────────────────────────────────────────
# SYSTEM SLIDES
# ─────────────────────────────────────────────

def slide_system(prs, d, sys_data, idx, total, slide_number):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_textbox(slide, LM, Inches(0.18), Inches(5), Inches(0.3),
                f"SYSTEM {idx} OF {total}", "Arial", 11, bold=True, color=C_PURPLE)

    name_len = len(sys_data["name"])
    title_size = 34 if name_len <= 20 else (28 if name_len <= 30 else 22)
    add_title(slide, sys_data["name"], top=Inches(0.55), font_size=title_size,
              width=Inches(7.6), height=Inches(0.65))

    badge = add_rounded_rect(slide, LM, Inches(1.28), Inches(2.6), Inches(0.36),
                              fill_color=C_DARK_TEXT, radius_pt=4)
    set_shape_text(badge, sys_data["label"], "Arial", 10, bold=True,
                   color=C_WHITE, align=PP_ALIGN.CENTER)
    badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, LM, Inches(1.78), Inches(7.5), Inches(0.95),
                sys_data["description"], "Calibri", 13, color=C_DARK_TEXT)

    add_textbox(slide, LM, Inches(2.82), Inches(4), Inches(0.25),
                "HOW IT WORKS", "Arial", 10, bold=True, color=C_PURPLE)

    steps = sys_data["steps"]
    step_w = (Inches(7.5) - Inches(0.08) * 4) / 5
    # Made step cards taller so 3-line step text fits without overflow.
    step_h = Inches(1.3)
    step_top = Inches(3.12)

    for i, step in enumerate(steps):
        cx = LM + i * (step_w + Inches(0.08))
        add_rounded_rect(slide, cx, step_top, step_w, step_h,
                         fill_color=C_HIGHLIGHT, radius_pt=4)
        add_circle(slide, cx + step_w / 2, step_top + Inches(0.25),
                   Inches(0.34), C_PURPLE, str(i + 1), font_size=12)
        # Step text occupies the BOTTOM ~0.8" of the card, below the circle
        add_vcenter_text(slide, cx + Inches(0.06), step_top + Inches(0.48),
                         step_w - Inches(0.12), step_h - Inches(0.54),
                         step, "Calibri", 9, color=C_DARK_TEXT)

    # AUDIT FINDING — shifted down so it doesn't collide with tall step cards
    af_top = step_top + step_h + Inches(0.15)   # ≈ 4.57"
    af_h = Inches(1.0)
    add_rect(slide, LM, af_top, Inches(7.5), af_h,
             fill_color=RGBColor(0xFF, 0xF0, 0xF0))
    add_rect(slide, LM, af_top, Inches(0.05), af_h, fill_color=C_RED)
    add_textbox(slide, LM + Inches(0.15), af_top + Inches(0.06),
                Inches(3), Inches(0.25),
                "AUDIT FINDING", "Arial", 10, bold=True, color=C_RED)
    add_textbox(slide, LM + Inches(0.15), af_top + Inches(0.32),
                Inches(7.2), Inches(0.65),
                sys_data["audit_finding"], "Calibri", 11, color=C_DARK_TEXT)

    roi_top = af_top + af_h + Inches(0.1)
    add_rounded_rect(slide, LM, roi_top, Inches(7.5), Inches(0.58),
                     fill_color=C_DARK_TEXT, radius_pt=4)
    add_textbox(slide, LM + Inches(0.2), roi_top + Inches(0.18),
                Inches(4.5), Inches(0.28),
                "PROJECTED ROI", "Arial", 10, bold=True, color=C_MUTED)
    add_textbox(slide, LM + Inches(4.0), roi_top + Inches(0.12),
                Inches(3.3), Inches(0.38),
                f"{sys_data['roi']} return on investment",
                "Calibri", 16, bold=True, color=C_GREEN, align=PP_ALIGN.RIGHT)

    # RIGHT COLUMN — 5 equal cards.
    # Key fix: label is at the TOP of the card with a reserved height, and
    # the value lives BELOW it in its own region. Previously labels sat at
    # 0.1" and the value at 0.38" within a 0.7" card, which caused the big
    # number to spill past the card top for the first card.
    rc_x = Inches(8.35)
    rc_w = W - rc_x - RM
    rc_top = Inches(0.18)
    rc_bottom = Inches(6.88)
    n_cards = 5
    gap = Inches(0.1)
    card_h = (rc_bottom - rc_top - gap * (n_cards - 1)) / n_cards

    cards_rc = [
        (sys_data["metric_label"],  str(sys_data["metric_value"]), C_DARK_TEXT, 32, C_CARD_BG),
        ("Monthly Revenue",          sys_data["monthly_revenue"],   C_GREEN,     26, C_CARD_BG),
        ("Monthly Cost",             sys_data["monthly_cost"],       C_AMBER,     26, C_CARD_BG),
        ("Setup Fee",                sys_data["setup_fee"],          C_DARK_TEXT, 26, C_CARD_BG),
        ("ROI",                      sys_data["roi"],                C_GREEN,     26, C_DARK_TEXT),
    ]

    for i, (label, value, val_color, val_size, bg) in enumerate(cards_rc):
        cy = rc_top + i * (card_h + gap)
        add_rounded_rect(slide, rc_x, cy, rc_w, card_h,
                         fill_color=bg, radius_pt=4,
                         line_color=None if bg == C_DARK_TEXT else RGBColor(0xE0, 0xE0, 0xE0))
        label_color = C_MUTED if bg == C_CARD_BG else RGBColor(0x99, 0x99, 0xBB)

        # Reserve a fixed label zone at the top of each card, then the value
        # box starts well below it. This prevents any label/value collision.
        label_h = Inches(0.5)   # room for 2 lines at 10pt
        add_textbox(slide, rc_x + Inches(0.12), cy + Inches(0.08),
                    rc_w - Inches(0.24), label_h,
                    label, "Calibri", 10, color=label_color)
        value_box_top = cy + Inches(0.08) + label_h + Inches(0.02)
        value_box_h = card_h - (value_box_top - cy) - Inches(0.1)
        if value_box_h < Inches(0.3):
            value_box_h = Inches(0.3)
        add_textbox(slide, rc_x + Inches(0.12), value_box_top,
                    rc_w - Inches(0.24), value_box_h,
                    value, "Georgia", val_size, bold=True, color=val_color)

    add_footer(slide, d["brand_name"], slide_number)


# ─────────────────────────────────────────────
# REVENUE SUMMARY
# ─────────────────────────────────────────────

def slide_revenue_summary(prs, d, slide_number):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_DARK_BG)
    add_accent_bar(slide, dark=True)

    add_textbox(slide, LM, Inches(0.3), Inches(8), Inches(0.28),
                "REVENUE RECOVERY SUMMARY", "Arial", 11, bold=True, color=C_PURPLE)

    add_textbox(slide, LM, Inches(0.65), Inches(8), Inches(1.5),
                d["total_monthly_revenue"], "Georgia", 80, bold=True, color=C_GREEN)
    add_textbox(slide, LM, Inches(2.1), Inches(8), Inches(0.4),
                f"additional revenue per month  \u00b7  {d['annual_recovery']} per year",
                "Calibri", 16, color=C_MUTED)

    # 4 KPI boxes — widened + auto-shrinking values so long strings fit.
    kpis = [
        ("Monthly Investment", d["total_monthly_cost"],     C_WHITE),
        ("Net Monthly Gain",   d["net_monthly_gain"],       C_GREEN),
        ("Annual Recovery",    d["annual_recovery"],        C_GREEN),
        ("Overall ROI",        d["overall_roi"] + " return", C_PURPLE),
    ]
    # Use full content width for the row
    kpi_row_w = W - LM - RM
    kpi_gap = Inches(0.14)
    kpi_w = (kpi_row_w - kpi_gap * 3) / 4
    kpi_h = Inches(1.0)
    kpi_top = Inches(2.62)

    for i, (label, value, color) in enumerate(kpis):
        cx = LM + i * (kpi_w + kpi_gap)
        add_rounded_rect(slide, cx, kpi_top, kpi_w, kpi_h,
                         fill_color=C_DARK_CARD, radius_pt=4)
        add_textbox(slide, cx + Inches(0.12), kpi_top + Inches(0.06),
                    kpi_w - Inches(0.24), Inches(0.28),
                    label, "Calibri", 10, color=C_MUTED)
        # Shrink value aggressively so long strings ("$97,200", "4.4x return")
        # fit inside a ~2.9" KPI card without clipping.
        vlen = len(value)
        if vlen <= 6:
            v_size = 26
        elif vlen <= 8:
            v_size = 22
        else:
            v_size = 20
        add_textbox(slide, cx + Inches(0.12), kpi_top + Inches(0.35),
                    kpi_w - Inches(0.24), Inches(0.58),
                    value, "Georgia", v_size, bold=True, color=color)

    systems = d["systems"]
    table_top = Inches(3.78)
    col_widths = [Inches(4.0), Inches(2.0), Inches(1.6), Inches(1.6), Inches(1.6)]
    headers = ["AI System", "Revenue/mo", "Cost/mo", "Setup", "ROI"]
    row_h = Inches(0.46)

    x = LM
    for j, (header, cw) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, table_top, cw, row_h, fill_color=C_PURPLE)
        add_textbox(slide, x + Inches(0.08), table_top + Inches(0.1),
                    cw - Inches(0.16), Inches(0.26),
                    header, "Calibri", 12, bold=True, color=C_WHITE)
        x += cw

    for r, sys in enumerate(systems):
        row_y = table_top + (r + 1) * row_h
        bg = C_DARK_CARD if r % 2 == 0 else RGBColor(0x24, 0x22, 0x3A)
        row_data = [
            (sys["name"],           C_WHITE,  False),
            (sys["monthly_revenue"], C_GREEN, False),
            (sys["monthly_cost"],    C_WHITE, False),
            (sys["setup_fee"],       C_WHITE, False),
            (sys["roi"],             C_PURPLE, False),
        ]
        x = LM
        for j, ((text, color, bold), cw) in enumerate(zip(row_data, col_widths)):
            add_rect(slide, x, row_y, cw, row_h, fill_color=bg)
            add_textbox(slide, x + Inches(0.08), row_y + Inches(0.1),
                        cw - Inches(0.16), Inches(0.26),
                        text, "Calibri", 12, bold=bold, color=color)
            x += cw

    total_y = table_top + (len(systems) + 1) * row_h
    total_data = [
        ("TOTAL",                   C_WHITE,  True),
        (d["total_monthly_revenue"] + "/mo", C_GREEN, True),
        (d["total_monthly_cost"] + "/mo",    C_WHITE, True),
        (d["total_setup_fees"],     C_WHITE,  True),
        (d["overall_roi"],          C_PURPLE, True),
    ]
    x = LM
    for (text, color, bold), cw in zip(total_data, col_widths):
        add_rect(slide, x, total_y, cw, row_h, fill_color=C_DARK_CARD)
        add_textbox(slide, x + Inches(0.08), total_y + Inches(0.1),
                    cw - Inches(0.16), Inches(0.26),
                    text, "Calibri", 12, bold=bold, color=color)
        x += cw

    add_footer(slide, d["brand_name"], slide_number)


# ─────────────────────────────────────────────
# TRANSFORMATION
# ─────────────────────────────────────────────

def slide_transformation(prs, d, slide_number):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_section_label(slide, "THE TRANSFORMATION")
    add_title(slide, "Before vs. After AI Systems", top=TITLE_TOP, font_size=32, auto_shrink=True)

    rows = d["transformation_rows"]
    table_top = CONTENT_TOP
    row_h = Inches(0.5)
    label_w = Inches(4.5)
    val_w = Inches(3.8)
    header_h = Inches(0.45)

    add_rect(slide, LM, table_top, label_w + val_w * 2 + Inches(0.04),
             header_h, fill_color=C_DARK_TEXT)
    add_textbox(slide, LM + label_w + Inches(0.1), table_top + Inches(0.08),
                val_w - Inches(0.2), Inches(0.28),
                "CURRENT STATE", "Arial", 10, bold=True, color=C_RED,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, LM + label_w + val_w + Inches(0.14), table_top + Inches(0.08),
                val_w - Inches(0.2), Inches(0.28),
                "WITH AI SYSTEMS", "Arial", 10, bold=True, color=C_GREEN,
                align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        ry = table_top + header_h + i * row_h
        bg = C_HIGHLIGHT if i % 2 == 0 else C_WHITE
        add_rect(slide, LM, ry, label_w + val_w * 2 + Inches(0.04), row_h,
                 fill_color=bg)
        add_textbox(slide, LM + Inches(0.12), ry + Inches(0.12),
                    label_w - Inches(0.24), Inches(0.26),
                    row["label"], "Calibri", 14, color=C_DARK_TEXT)
        add_textbox(slide, LM + label_w + Inches(0.1), ry + Inches(0.1),
                    val_w - Inches(0.2), Inches(0.3),
                    row["current"], "Calibri", 14, bold=True, color=C_RED,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, LM + label_w + val_w + Inches(0.14), ry + Inches(0.1),
                    val_w - Inches(0.2), Inches(0.3),
                    row["with_ai"], "Calibri", 14, bold=True, color=C_GREEN,
                    align=PP_ALIGN.CENTER)

    add_footer(slide, d["brand_name"], slide_number)


# ─────────────────────────────────────────────
# GUARANTEE
# ─────────────────────────────────────────────

def slide_guarantee(prs, d, slide_number):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_DARK_BG)

    add_rect(slide, Inches(0.5), Inches(0.25), W - Inches(1.0), Inches(0.06),
             fill_color=C_TEAL)
    add_rect(slide, Inches(0.5), H - Inches(0.31), W - Inches(1.0), Inches(0.06),
             fill_color=C_TEAL)

    add_rounded_rect(slide, Inches(0.5), Inches(0.35),
                     W - Inches(1.0), H - Inches(0.7),
                     fill_color=C_DARK_CARD, radius_pt=6)

    add_textbox(slide, 0, Inches(0.52), W, Inches(0.3),
                "OUR GUARANTEE", "Arial", 11, bold=True, color=C_PURPLE,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(0.88), W - Inches(2.0), Inches(1.65),
                d["guarantee_headline"],
                "Georgia", 40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(2.7), W - Inches(2.0), Inches(0.48),
                "Or we work for free until you do.",
                "Georgia", 22, italic=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    # Math box — 3 tightly-packed rows. Rewrote so label + value appear as
    # two runs on a single paragraph (value coloured green), instead of the
    # earlier version where values were plain white inside one combined run.
    math_y = Inches(3.45)
    math_w = W - Inches(3.0)
    math_h = Inches(3.1)
    add_rounded_rect(slide, Inches(1.5), math_y, math_w, math_h,
                     fill_color=RGBColor(0x28, 0x26, 0x42), radius_pt=4)

    add_textbox(slide, Inches(1.7), math_y + Inches(0.18), math_w - Inches(0.4), Inches(0.3),
                "THE MATH", "Arial", 11, bold=True, color=C_MUTED,
                align=PP_ALIGN.CENTER)

    # Three rows with label (white) + value (green) on the same line.
    math_items = [
        (f"Your 12-month investment:", d["guarantee_investment"], C_WHITE, C_GREEN),
        (f"Projected revenue in first 90 days:", d["guarantee_90day"], C_WHITE, C_GREEN),
        (f"That's {d['guarantee_surplus']} surplus \u2014 before you've even used 3 months.",
         "", C_GREEN, C_GREEN),
    ]
    row_start = math_y + Inches(0.7)
    row_spacing = Inches(0.75)
    for idx, (label, value, lcolor, vcolor) in enumerate(math_items):
        ry = row_start + idx * row_spacing
        # Compose label + value as two-run paragraph for correct colour
        txBox = slide.shapes.add_textbox(Inches(1.7), ry,
                                         math_w - Inches(0.4), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run1 = p.add_run()
        run1.text = label
        run1.font.name = "Calibri"
        run1.font.size = Pt(18)
        run1.font.color.rgb = lcolor
        if value:
            run2 = p.add_run()
            run2.text = "  " + value
            run2.font.name = "Georgia"
            run2.font.size = Pt(22)
            run2.font.bold = True
            run2.font.color.rgb = vcolor
        # Zero padding
        bodyPr = tf._txBody.find(_qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('lIns', '0'); bodyPr.set('rIns', '0')
            bodyPr.set('tIns', '0'); bodyPr.set('bIns', '0')


# ─────────────────────────────────────────────
# IMPLEMENTATION
# ─────────────────────────────────────────────

def slide_implementation(prs, d, slide_number):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_LIGHT_BG)
    add_accent_bar(slide, dark=False)

    add_section_label(slide, "IMPLEMENTATION")
    add_title(slide, "Live in 14 days. Results in 30.", top=TITLE_TOP, font_size=36)

    phases = d["implementation_phases"]
    col_w = Inches(3.85)
    col_h = Inches(4.2)
    col_top = CONTENT_TOP
    gap = Inches(0.22)

    phase_colors = [C_PURPLE, C_GREEN, RGBColor(0x00, 0x88, 0xCC)]

    for i, phase in enumerate(phases):
        cx = LM + i * (col_w + gap)
        add_rounded_rect(slide, cx, col_top, col_w, col_h,
                         fill_color=C_HIGHLIGHT, radius_pt=6)
        # Period label
        add_textbox(slide, cx + Inches(0.15), col_top + Inches(0.12),
                    col_w - Inches(0.3), Inches(0.25),
                    phase["period"], "Arial", 10, bold=True,
                    color=phase_colors[i])
        # Title — shrink font if title is long; reserve TWO lines of room
        # so the first item never collides with a wrapped title.
        title_text = phase["title"]
        title_fs = 17 if len(title_text) <= 26 else (15 if len(title_text) <= 34 else 13)
        title_h = Inches(0.85)   # enough for 2 lines at 15pt
        add_textbox(slide, cx + Inches(0.15), col_top + Inches(0.4),
                    col_w - Inches(0.3), title_h,
                    title_text, "Georgia", title_fs, bold=True, color=C_DARK_TEXT)
        # Items start below the reserved title area (was 0.9 — now 1.35)
        items_top = col_top + Inches(1.35)
        for j, item in enumerate(phase["items"]):
            add_textbox(slide, cx + Inches(0.15), items_top + j * Inches(0.42),
                        col_w - Inches(0.3), Inches(0.38),
                        "\u2192  " + item, "Calibri", 12, color=C_DARK_TEXT)

    add_textbox(slide, LM, col_top + col_h + Inches(0.15), W - LM - RM, Inches(0.35),
                d["implementation_closing"], "Calibri", 13, italic=True,
                color=C_MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, d["brand_name"], slide_number)


# ─────────────────────────────────────────────
# NEXT STEPS
# ─────────────────────────────────────────────

def slide_next_steps(prs, d, slide_number):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_DARK_BG)
    add_accent_bar(slide, dark=True)

    add_textbox(slide, LM, Inches(0.32), Inches(5), Inches(0.28),
                "NEXT STEPS", "Arial", 11, bold=True, color=C_PURPLE)
    add_textbox(slide, LM, Inches(0.65), W - LM - RM, Inches(0.58),
                "Select your bots. We'll start building.",
                "Georgia", 32, bold=True, color=C_WHITE)

    steps = d["next_steps"]
    card_w = Inches(2.9)
    card_h = Inches(4.2)
    gap = Inches(0.24)
    top_y = Inches(1.45)

    for i, step in enumerate(steps):
        cx = LM + i * (card_w + gap)
        add_rounded_rect(slide, cx, top_y, card_w, card_h,
                         fill_color=C_DARK_CARD, radius_pt=6)
        add_textbox(slide, cx + Inches(0.18), top_y + Inches(0.15),
                    card_w - Inches(0.36), Inches(1.0),
                    str(i + 1), "Georgia", 60, bold=True, color=C_PURPLE)
        add_textbox(slide, cx + Inches(0.18), top_y + Inches(1.18),
                    card_w - Inches(0.36), Inches(0.46),
                    step["title"], "Calibri", 14, bold=True, color=C_WHITE)
        add_textbox(slide, cx + Inches(0.18), top_y + Inches(1.68),
                    card_w - Inches(0.36), Inches(1.8),
                    step["description"], "Calibri", 12, color=C_MUTED)

    add_textbox(slide, LM, Inches(5.85), W - LM - RM, Inches(0.32),
                f"{d['cta_url']}  \u00b7  {d['cta_label']}",
                "Calibri", 13, color=C_MUTED, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# CLOSING SLIDE
# ─────────────────────────────────────────────

def slide_closing(prs, d, slide_number):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, W, H, fill_color=C_DARK_BG)
    add_accent_bar(slide, dark=True)

    # Urgency line — uses safe multi-line rendering; width reduced to prevent
    # single giant run. Height is generous so 3-4 wrapped lines fit.
    add_textbox(slide, LM, Inches(1.6), W - LM - RM, Inches(2.6),
                d["closing_urgency_line"],
                "Georgia", 28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # CTA — pushed to its own y-band well below the urgency block so they
    # cannot overlap no matter how long the urgency copy is.
    add_textbox(slide, LM, Inches(4.6), W - LM - RM, Inches(1.0),
                d["closing_cta"],
                "Georgia", 44, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    add_textbox(slide, LM, H - Inches(0.75), W - LM - RM, Inches(0.3),
                d["brand_line"],
                "Calibri", 14, color=C_MUTED, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# MAIN GENERATE FUNCTION
# ─────────────────────────────────────────────

def generate(data, output_path="output.pptx"):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_01_cover(prs, data)
    slide_02_agenda(prs, data)
    slide_03_part_one(prs, data)
    slide_04_journey(prs, data)
    slide_05_stats(prs, data)
    slide_06_lifecycle(prs, data)
    slide_07_funnel(prs, data)
    slide_08_part_two(prs, data)
    slide_09_numbers(prs, data)
    slide_11_score(prs, data)

    current_slide = 11
    for dive in data["deep_dives"]:
        slide_deep_dive(prs, data, dive, current_slide)
        current_slide += 1

    slide_15_bottom_line(prs, data)
    current_slide += 1

    slide_16_part_three(prs, data)
    current_slide += 1

    slide_17_how_it_works(prs, data)
    current_slide += 1

    systems = data["systems"]
    total_systems = len(systems)
    for i, sys in enumerate(systems):
        slide_system(prs, data, sys, i + 1, total_systems, current_slide)
        current_slide += 1

    slide_revenue_summary(prs, data, current_slide); current_slide += 1
    slide_transformation(prs, data, current_slide);  current_slide += 1
    slide_guarantee(prs, data, current_slide);        current_slide += 1
    slide_implementation(prs, data, current_slide);   current_slide += 1
    slide_next_steps(prs, data, current_slide);       current_slide += 1
    slide_closing(prs, data, current_slide)

    for slide in prs.slides:
        add_logo_watermark(slide)

    prs.save(output_path)
    print(f"\u2713 Saved: {output_path}  ({current_slide} slides)")
    return output_path


# ─────────────────────────────────────────────
# TEST DATA — matches the deployed RE Business LLC deck
# ─────────────────────────────────────────────

TEST_DATA = {
    "brand_name":    "justin.babcock98@gmail.com",
    "client_name":   "RE Business LLC",
    "prepared_for":  "Justin Babcock",
    "date":          "21 April 2026",
    "cta_url":       "agencypartnerhub.lovable.app/p/justin-babcock98-gmail-com",
    "cta_label":     "Your personalised bot selector is ready",
    "brand_line":    "justin.babcock98@gmail.com  \u00b7  agencypartnerhub.lovable.app",

    "term_client":       "lead",
    "term_client_pl":    "leads",
    "term_engagement":   "appraisal",
    "term_ongoing":      "listing",
    "term_completing":   "settling the sale",
    "term_practitioner": "agent",
    "term_business":     "agency",
    "bottom_line_biz_word": "agency",

    "agenda_01_desc": "The buyer & vendor journey and where deals slip through",

    "part_one_title":    "Your Industry &\nThe Client Journey",
    "part_one_subtitle": "Understanding where revenue leaks in a real estate agency",

    "journey_label":   "THE BUYER & VENDOR JOURNEY",
    "journey_title":   "How a real estate lead finds you \u2014 and where they drop off",
    "journey_stages": [
        {"name": "AWARENESS",    "desc": "Portals, signs,\nreferrals",               "pct_lost": None,  "reason": ""},
        {"name": "ENQUIRY",      "desc": "Portal lead, web form,\nappraisal request","pct_lost": "48%", "reason": "No response in first 5 min"},
        {"name": "RESPONSE",     "desc": "Agent calls back",                          "pct_lost": "30%", "reason": "Never followed up again"},
        {"name": "APPRAISAL",    "desc": "In-home appraisal\nor inspection",          "pct_lost": None,  "reason": ""},
        {"name": "LISTING",      "desc": "Signs authority\nor buyer agreement",       "pct_lost": "50%", "reason": "Listed with faster agent"},
        {"name": "LIFETIME",     "desc": "Referrals, repeat,\nfuture transactions",   "pct_lost": "80%", "reason": "No post-settlement contact"},
    ],
    "journey_closing": "Every red zone is a commission you generated interest for \u2014 but never collected.",

    "stat_section_label": "REAL ESTATE INDUSTRY DATA",
    "industry_stats": [
        {"number": "48%",   "label": "of real estate leads are never contacted at all \u2014 most go to whichever agent calls first, not the one who was emailed first", "source": "NAR Real Estate Lead Response Study, 2025"},
        {"number": "5min",  "label": "window before a portal lead's interest drops by 80% \u2014 most agents take 30+ minutes to respond and lose the listing",            "source": "Harvard Business Review Lead Response Research"},
        {"number": "12%",   "label": "average conversion from enquiry to signed authority \u2014 top agencies hit 30%+ by fixing speed, follow-up, and nurture",           "source": "REB Agency Benchmarks, 2025"},
        {"number": "8x",    "label": "more likely to win a listing when you're the first agent to respond vs the third \u2014 speed-to-lead is the single biggest predictor", "source": "Zillow Premier Agent Data, 2024"},
    ],

    "lifecycle_title":   "A lead isn't a one-time transaction \u2014 it's a lifecycle",
    "lifecycle_intro":   "Most agencies only focus on winning the next listing. But sustainable GCI comes from managing the full client lifecycle \u2014 and most offices have no systems for the stages marked below.",
    "lifecycle_stages": [
        {"name": "ATTRACT", "desc": "Portal spend, signs,\nfarm area marketing",  "status": "WORKING"},
        {"name": "CAPTURE", "desc": "Answer calls,\nrespond to portal leads",      "status": "WORKING"},
        {"name": "NURTURE", "desc": "Follow up leads\nwho didn't list/buy yet",    "status": "NO SYSTEM"},
        {"name": "CONVERT", "desc": "Book appraisal,\nwin the listing",            "status": "MANUAL"},
        {"name": "RETAIN",  "desc": "Keep vendors/buyers\nupdated through deal",   "status": "NO SYSTEM"},
        {"name": "GROW",    "desc": "Reviews, referrals,\nrepeat transactions",    "status": "NO SYSTEM"},
    ],
    "lifecycle_summary": "4 of 6 lifecycle stages have no system. You're paying to attract leads, then losing them to faster-responding agents.",

    "funnel_rows": [
        {"label": "Monthly Enquiries (portal + web + signs)", "value": 15, "lost": "-5 lost"},
        {"label": "Answered / Responded To (within 1hr)",     "value": 10, "lost": "-4 lost"},
        {"label": "Actually Followed Up (if didn't list)",    "value": 6,  "lost": "-5 lost"},
        {"label": "Booked Appraisal / Inspection",            "value": 1,  "lost": ""},
        {"label": "Signed Authority / Agreement",             "value": 1,  "lost": None},
    ],
    "funnel_insight": "Of 15 monthly enquiries, only 10 get a response, 6 see follow-up, and just 1 books an appraisal \u2014 a 93% drop-off between the top of the funnel and a signed listing.",

    "number_cards": [
        {"value": "15",      "label": "Monthly Enquiries",                        "subnote": "Portal, web, signs, referrals"},
        {"value": "~2",      "label": "Missed / Unanswered\nCalls per Month",    "subnote": "No after-hours coverage"},
        {"value": "7%",      "label": "Close Rate\n(enquiry \u2192 authority)",  "subnote": "From enquiry to signed authority"},
        {"value": "$150",    "label": "Average First\nAppraisal Value",          "subnote": "Time-cost of an appraisal"},
        {"value": "$12,000", "label": "Avg Commission per Sale",                 "subnote": "GCI per transaction"},
        {"value": "99%",     "label": "Listing\nConversion Rate",                "subnote": "From appraisal to authority"},
    ],

    "audit_tests": [
        {"name": "Phone Call Test (business hours)",  "finding": "Called at 11am Wednesday. Rang 7 times.\nVoicemail \u2014 callback took 2+ hours"},
        {"name": "Phone Call Test (after hours)",     "finding": "Called at 6:30pm Thursday.\nVoicemail \u2014 no callback until next morning"},
        {"name": "Web Enquiry Test",                  "finding": "Submitted web form at 10am Tuesday.\nEmail reply 4 hours later, no phone call"},
        {"name": "Portal Lead Flow Test",             "finding": "Tested portal intake flow end-to-end.\nWorks but no confirmation SMS sent"},
        {"name": "Follow-Up Persistence",             "finding": "After web form, we didn't respond.\nZero follow-up attempts received"},
        {"name": "Google Review Audit",               "finding": "Compared profile vs top 5 local agents.\nFewer reviews than top local agencies"},
    ],

    "overall_score": 38,
    "score_status_label": "CRITICAL \u2014 IMMEDIATE\nACTION NEEDED",
    "dimensions": [
        {"name": "Speed to Lead",         "weight": "25%", "score": 46},
        {"name": "Follow-Up Systems",     "weight": "20%", "score": 50},
        {"name": "Pipeline Visibility",   "weight": "15%", "score": 45},
        {"name": "Vendor & Buyer Comm.",  "weight": "15%", "score": 23},
        {"name": "Post-Settlement Nurture","weight": "15%", "score": 20},
        {"name": "Automation Maturity",   "weight": "10%", "score": 35},
    ],
    "key_findings": [
        "Overall sales maturity scores 38 out of 100, with Post-Settlement Nurture bottoming out at 20.",
        "$5,922 in revenue leaks out of the agency every month, or $71,064 annually \u2014 47% of what you should be earning.",
        "Only 1 of 15 monthly enquiries converts to a booked appraisal, a 93% funnel drop-off.",
        "Response time sits in the 5-30 minute window, well past the 5-minute threshold where lead conversion collapses.",
        "With a 7% close rate on $12,000 average deals, each recovered lead is worth real money \u2014 GoHighLevel is currently underused.",
    ],

    "deep_dives": [
        {
            "score": 23,
            "name": "Vendor & Buyer Communication",
            "subtitle": "Vendor & Buyer Communication",
            "sub_scores": [
                {"score": 28, "label": "Cadence Consistency",  "note": "Touchpoints are ad hoc, not scheduled."},
                {"score": 22, "label": "Channel Mix",          "note": "Email-heavy, minimal SMS or voice."},
                {"score": 20, "label": "Template Quality",     "note": "No branded, niche-specific templates in use."},
                {"score": 22, "label": "Response Handling",    "note": "Inbound vendor replies often sit over 24 hours."},
            ],
            "audit_finding":  "RE Business LLC relies on the agent's manual discipline to keep vendors and buyers informed between appraisal and settlement. With 15 enquiries monthly, the cracks show fast \u2014 leads feel ignored and walk to competitors who communicate proactively.",
            "revenue_impact": "Communication gaps are a core driver of the $5,922 monthly leak. Even recovering 2 of every 15 leads at a $12,000 deal value and 7% close rate meaningfully lifts monthly revenue above the current $10,500.",
        },
        {
            "score": 20,
            "name": "Post-Settlement Nurture",
            "subtitle": "Post-Settlement Nurture",
            "sub_scores": [
                {"score": 18, "label": "Referral Asks",          "note": "No systematic referral trigger after settling the sale."},
                {"score": 20, "label": "Review Capture",          "note": "Reviews collected sporadically, not automated."},
                {"score": 22, "label": "Anniversary Touches",    "note": "No 12-month check-in cadence in place."},
                {"score": 20, "label": "Reactivation",            "note": "Past clients never re-engaged for new listings."},
            ],
            "audit_finding":  "Scoring 20 out of 100, Post-Settlement Nurture is the single biggest weakness at RE Business LLC. Once a sale settles, the lead relationship effectively ends \u2014 no reviews, no referrals, no reactivation pipeline feeding future listings.",
            "revenue_impact": "A Review Bot and reactivation layer alone is modeled at $2,400 in monthly revenue, a direct recovery of a large slice of the $5,922 monthly leak.",
        },
        {
            "score": 35,
            "name": "Automation Maturity",
            "subtitle": "Automation Maturity",
            "sub_scores": [
                {"score": 38, "label": "GoHighLevel Utilization", "note": "Platform in place, workflows underbuilt."},
                {"score": 32, "label": "Lead Routing",             "note": "No instant routing to agent on hot enquiries."},
                {"score": 35, "label": "Sequence Depth",           "note": "Only 3-5 touches per lead, no branching logic."},
                {"score": 35, "label": "Reporting",                "note": "No automated leak or conversion dashboards."},
            ],
            "audit_finding":  "GoHighLevel is already paid for at RE Business LLC but operating at a fraction of its capability. Without deeper automation, the agency cannot scale beyond its current under-20 monthly lead volume.",
            "revenue_impact": "Activating proper automation unlocks the modeled $8,100 net monthly gain and keeps the $71,064 annual leak from repeating next year.",
        },
    ],

    "annual_leak":          "$71,064",
    "bottom_line_context":  "This breakdown accounts for the $5,922 in monthly revenue leaking out of RE Business LLC across slow lead response, weak follow-up, and zero post-settlement nurture. Annualized, that is $71,064 walking out the door every year.",
    "closing_line":         "The question isn't whether you can afford AI systems.\nIt's whether you can afford not to have them.",

    "before_dropoff_pct":     "47% drop off",
    "after_step_5":           "Vendor update bot\nkeeps listings active",
    "how_it_works_closing":   "Every stage is automated. Your team focuses on winning listings \u2014 AI handles the rest.",

    "systems": [
        {
            "name":            "AI Receptionist",
            "label":           "PHASE 1 \u2014 CORE",
            "metric_label":    "minute response\nwindow target",
            "metric_value":    "5",
            "monthly_revenue": "$4,500",
            "monthly_cost":    "$1,000",
            "setup_fee":       "$2,000",
            "roi":             "4.5x",
            "description":     "An always-on AI receptionist that answers inbound enquiries within seconds, qualifies the lead, and books appraisals straight into the agent's calendar. Sits on top of GoHighLevel so nothing is replaced \u2014 only accelerated.",
            "steps":           [
                "Capture inbound\nenquiry across\nweb, SMS, voice",
                "Qualify lead with\nniche-specific\nquestions",
                "Route hot leads\nto agent instantly",
                "Book appraisal\ndirectly into\ncalendar",
                "Log everything\ninto GoHighLevel",
            ],
            "audit_finding":   "With response time sitting in the 5-30 minute band and 2 missed calls per month, RE Business LLC is losing leads before the agent even sees them. An AI receptionist closes that gap and is modeled at $4,500 monthly revenue recovery.",
        },
        {
            "name":            "Lead Nurture Bot",
            "label":           "PHASE 2 \u2014 NURTURE",
            "metric_label":    "automated touches\nper lead",
            "metric_value":    "10+",
            "monthly_revenue": "$3,600",
            "monthly_cost":    "$600",
            "setup_fee":       "$1,500",
            "roi":             "6.0x",
            "description":     "A multi-channel nurture bot that runs structured SMS, email, and voice sequences from first enquiry through settling the sale. Replaces the current 3-5 manual touches with a consistent cadence that never forgets a lead.",
            "steps":           [
                "Trigger sequence\non enquiry capture",
                "Multi-channel touches\nacross SMS and email",
                "Branch logic based\non lead behavior",
                "Re-engage cold leads\nat 30, 60, 90 days",
                "Hand warm leads\nback to the agent",
            ],
            "audit_finding":   "Only 6 of 15 monthly enquiries currently receive follow-up at RE Business LLC. A Lead Nurture Bot extends coverage to 100% of the pipeline and is modeled at $3,600 monthly revenue.",
        },
        {
            "name":            "Review Bot + Client Reactivation",
            "label":           "ADD-ON \u2014 GROWTH",
            "metric_label":    "reactivation touches\nper past client",
            "metric_value":    "4",
            "monthly_revenue": "$2,400",
            "monthly_cost":    "$800",
            "setup_fee":       "$2,000",
            "roi":             "3.0x",
            "description":     "Automated review capture the moment a sale settles, plus a reactivation engine that re-engages past clients for referrals, anniversary check-ins, and new listing opportunities. This is where Post-Settlement Nurture moves from 20 to best-in-class.",
            "steps":           [
                "Trigger review\nrequest at settlement",
                "Route 5-star reviews\nto Google and Facebook",
                "Schedule anniversary\ncheck-in touches",
                "Re-engage past clients\nfor referrals",
                "Surface reactivation\nleads to the agent",
            ],
            "audit_finding":   "With zero post-settlement workflow in place at RE Business LLC, every past client is a dormant asset. This system is modeled at $2,400 monthly revenue, directly attacking the weakest score on the audit.",
        },
    ],

    "total_monthly_revenue": "$10,500",
    "total_monthly_cost":    "$2,400",
    "total_setup_fees":      "$5,500",
    "net_monthly_gain":      "$8,100",
    "annual_recovery":       "$97,200",
    "overall_roi":           "4.4x",

    "guarantee_headline":   "Make your full 12-month investment back in 90 days.",
    "guarantee_investment": "$28,800",
    "guarantee_90day":      "$31,500",
    "guarantee_surplus":    "$2,700",

    "transformation_rows": [
        {"label": "Response Time",         "current": "5-30 minutes",           "with_ai": "Under 60 seconds"},
        {"label": "After-Hours Coverage",  "current": "Auto-reply only",         "with_ai": "24/7 qualification and booking"},
        {"label": "Follow-Up Touches",     "current": "3-5 per lead",            "with_ai": "10+ multi-channel per lead"},
        {"label": "Funnel Conversion",     "current": "1 of 15 enquiries booked","with_ai": "3-4 of 15 enquiries booked"},
        {"label": "Post-Settlement Nurture","current": "None",                   "with_ai": "Automated reviews and reactivation"},
        {"label": "Monthly Revenue Leak",  "current": "$5,922",                  "with_ai": "Recovered into pipeline"},
        {"label": "Net Monthly Gain",      "current": "$0",                      "with_ai": "$8,100"},
    ],

    "implementation_phases": [
        {
            "period": "WEEK 1\u20132",
            "title":  "Deploy AI Receptionist",
            "items":  [
                "Integrate with GoHighLevel",
                "Configure qualification script",
                "Connect calendar booking",
                "Go live on inbound channels",
            ],
        },
        {
            "period": "WEEK 3\u20134",
            "title":  "Launch Lead Nurture Bot",
            "items":  [
                "Build multi-channel sequences",
                "Set branching logic by behavior",
                "Migrate existing leads into cadence",
                "Activate cold-lead re-engagement",
            ],
        },
        {
            "period": "MONTH 2\u20133",
            "title":  "Add Review & Reactivation Layer",
            "items":  [
                "Trigger post-settlement reviews",
                "Build anniversary touch cadence",
                "Launch past-client reactivation",
                "Report on recovered revenue",
            ],
        },
    ],
    "implementation_closing": "Zero disruption to your team. We handle the entire build, CRM integration, and testing.",

    "next_steps": [
        {"title": "Select your systems",            "description": "Use the link we'll send you to choose which bots you want and customise your numbers."},
        {"title": "We build everything in 14 days", "description": "Complete setup, CRM integration, and testing \u2014 zero work required from your team."},
        {"title": "Revenue starts recovering",      "description": "Most clients see measurable results within the first 30 days."},
        {"title": "90-day guarantee kicks in",      "description": "If you haven't made your 12-month investment back in 90 days, we work for free until you do."},
    ],

    "closing_urgency_line": "Every month without these systems costs RE Business LLC another $8,100 in net gain.\nThat is not a projection, that is the math on leads you are already paying to generate.",
    "closing_cta":          "Let's plug the leaks.",
}


if __name__ == "__main__":
    generate(TEST_DATA, "RE_Business_LLC_Revenue_Audit_FIXED.pptx")
